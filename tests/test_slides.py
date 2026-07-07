"""Slide generation: ported behavior plus the audited fixes, verified by
re-opening the generated files with python-pptx."""

import io

from PIL import Image
from pptx import Presentation

from nonnewtonian.parser import parse_file
from nonnewtonian.photos import extract_pptx_images, store_bytes
from nonnewtonian.slides import DeckChapter, add_entry_slides, build_deck, build_entry_slide, format_notes

from conftest import FIXTURES, SCIENTISTS


def _local_images(pptx_name, tmp_path, count=None):
    images = extract_pptx_images(FIXTURES / pptx_name)
    if count is not None:
        images = images[:count]
    return [store_bytes(data, tmp_path).path for data, _ in images]


def _reopen(prs) -> Presentation:
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return Presentation(buffer)


def test_one_slide_per_photo_with_notes(tmp_path):
    """Chien-Shiung Wu has two photos -> a 2-slide file, writeup in notes."""
    entry = parse_file(SCIENTISTS / "Chien-ShiungWu.txt")
    images = _local_images("Chien-Shiung Wu.pptx", tmp_path)
    prs = _reopen(build_entry_slide(entry, images))
    assert len(list(prs.slides)) == 2
    for slide in prs.slides:
        assert slide.shapes.title.text == "Chien-Shiung Wu"
        notes = slide.notes_slide.notes_text_frame.text
        assert "Chien-Shiung Wu" in notes
        assert entry.description[0][:40] in notes


def test_notes_are_formatted_not_raw_markup(tmp_path):
    """The original put the raw file text (headers and all) in notes."""
    entry = parse_file(SCIENTISTS / "EmmyNoether.txt")
    notes = format_notes(entry)
    assert "# Name" not in notes and "# Description" not in notes
    assert "Sources:" in notes


def test_photoless_entry_gets_a_text_slide():
    """Ursula Franklin has no photo; the original skipped her entirely."""
    entry = parse_file(SCIENTISTS / "UrsulaFranklin.txt")
    prs = _reopen(build_entry_slide(entry, []))
    slides = list(prs.slides)
    assert len(slides) == 1
    assert slides[0].shapes.title.text == "Ursula Franklin"


def test_image_aspect_fit_stays_on_slide(tmp_path):
    """A very wide image must be scaled to fit the slide width (the
    original's fixed 5-inch height overflowed on wide images)."""
    buffer = io.BytesIO()
    Image.new("RGB", (4000, 500), (10, 10, 200)).save(buffer, format="JPEG")
    path = store_bytes(buffer.getvalue(), tmp_path).path
    entry = parse_file(SCIENTISTS / "EmmyNoether.txt")
    prs = Presentation()
    add_entry_slides(prs, entry, [path])
    picture = next(
        shape for slide in prs.slides for shape in slide.shapes if shape.shape_type == 13
    )
    assert picture.left >= 0
    assert picture.left + picture.width <= prs.slide_width
    assert picture.top + picture.height <= prs.slide_height


def test_deck_is_chapter_ordered_with_dividers(tmp_path):
    noether = parse_file(SCIENTISTS / "EmmyNoether.txt")
    wu = parse_file(SCIENTISTS / "Chien-ShiungWu.txt")
    franklin = parse_file(SCIENTISTS / "UrsulaFranklin.txt")
    wu_images = _local_images("Chien-Shiung Wu.pptx", tmp_path, count=1)
    noether_images = _local_images("Emmy Noether.pptx", tmp_path, count=1)

    chapters = [
        DeckChapter(chapter=30, topics="Nuclear Physics", entries=[(wu, wu_images)]),
        DeckChapter(chapter=9, topics="Work and Energy", entries=[(noether, noether_images)]),
        DeckChapter(chapter=10, topics="Interactions", entries=[(noether, noether_images), (franklin, [])]),
    ]
    prs = _reopen(build_deck(chapters, deck_title="Test Deck"))
    titles = [slide.shapes.title.text for slide in prs.slides]
    assert titles[0] == "Test Deck"
    # Chapter order 9, 10, 30 regardless of input order; Noether under
    # both 9 and 10 by design; Franklin present despite no photo.
    assert titles[1].startswith("Chapter 9") and "Work and Energy" in titles[1]
    assert titles[2] == "Emmy Noether"
    assert titles[3].startswith("Chapter 10")
    assert "Emmy Noether" in titles[4:6] and "Ursula Franklin" in titles[4:6]
    assert titles[6].startswith("Chapter 30")
    assert titles[7] == "Chien-Shiung Wu"


def test_deck_dedupes_within_a_chapter(tmp_path):
    noether = parse_file(SCIENTISTS / "EmmyNoether.txt")
    images = _local_images("Emmy Noether.pptx", tmp_path, count=1)
    chapters = [
        DeckChapter(chapter=9, topics="", entries=[(noether, images), (noether, images)])
    ]
    prs = _reopen(build_deck(chapters))
    titles = [slide.shapes.title.text for slide in prs.slides]
    assert titles.count("Emmy Noether") == 1  # the original duplicated
