"""PPTX slide generation, ported from IntroductoryPhysics/makesyllabus.py
(add_scientist_slide / maketextbook) with the audited fixes:

- images come from LOCAL files, never the network at render time (one
  dead URL could abort a whole deck in the original);
- images are aspect-fit to the slide (the fixed 5-inch height could
  overflow the slide width);
- the notes pane gets the FORMATTED writeup (description paragraphs,
  sources, contributors, placements), not the raw file markup;
- entries without a photo get a text slide instead of being skipped
  (Ursula Franklin finally gets a slide);
- decks are chapter-ordered with a divider slide per chapter and are
  deduplicated per placement, instead of glob-order with duplicates.

Kept from the original on purpose: one slide per photo, so multi-photo
scientists (Chien-Shiung Wu has two) get multi-slide files.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Inches, Pt

from .parser import Entry

_TITLE_ONLY_LAYOUT = 5
_MARGIN_TOP = Inches(1.75)
_MARGIN_SIDE = Inches(0.5)
_MARGIN_BOTTOM = Inches(0.4)


@dataclass
class DeckChapter:
    """One chapter's worth of a deck: heading info + its entries."""

    chapter: int
    topics: str
    entries: list[tuple[Entry, list[Path]]] = field(default_factory=list)


def format_notes(entry: Entry, placements_text: list[str] | None = None) -> str:
    """The human-readable writeup that goes in the slide's notes pane."""
    parts: list[str] = [entry.name]
    placements = placements_text if placements_text is not None else entry.placements_raw
    if placements:
        parts.append("Placements:\n" + "\n".join(placements))
    if entry.description:
        parts.append("\n\n".join(entry.description))
    if entry.sources:
        parts.append("Sources:\n" + "\n".join(entry.sources))
    if entry.contributors:
        parts.append("Contributed by: " + ", ".join(entry.contributors))
    return "\n\n".join(parts)


def _fit_picture(slide_width: int, slide_height: int, image_path: Path):
    """Aspect-fit box (left, top, width, height) in EMU for the image."""
    with Image.open(image_path) as image:
        pixel_width, pixel_height = image.size
    available_width = slide_width - 2 * _MARGIN_SIDE
    available_height = slide_height - _MARGIN_TOP - _MARGIN_BOTTOM
    scale = min(available_width / pixel_width, available_height / pixel_height)
    width = int(pixel_width * scale)
    height = int(pixel_height * scale)
    left = (slide_width - width) // 2
    return Emu(left), Emu(_MARGIN_TOP), Emu(width), Emu(height)


def _add_title_slide(prs: Presentation, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[_TITLE_ONLY_LAYOUT])
    slide.shapes.title.text = title
    return slide

def add_entry_slides(
    prs: Presentation,
    entry: Entry,
    image_paths: list[Path],
    placements_text: list[str] | None = None,
) -> int:
    """Add this entry's slide(s) to a presentation.

    One slide per image; entries with no image get one text slide with
    the first description paragraph on the slide body.  Returns the
    number of slides added (always >= 1).
    """
    notes = format_notes(entry, placements_text)
    if not image_paths:
        slide = _add_title_slide(prs, entry.name)
        if entry.description:
            left, top = _MARGIN_SIDE, _MARGIN_TOP
            width = prs.slide_width - 2 * _MARGIN_SIDE
            height = prs.slide_height - _MARGIN_TOP - _MARGIN_BOTTOM
            box = slide.shapes.add_textbox(left, top, width, height)
            frame = box.text_frame
            frame.word_wrap = True
            frame.text = entry.description[0]
            frame.paragraphs[0].font.size = Pt(20)
        slide.notes_slide.notes_text_frame.text = notes
        return 1

    for image_path in image_paths:
        slide = _add_title_slide(prs, entry.name)
        left, top, width, height = _fit_picture(
            prs.slide_width, prs.slide_height, image_path
        )
        slide.shapes.add_picture(str(image_path), left, top, width, height)
        slide.notes_slide.notes_text_frame.text = notes
    return len(image_paths)


def build_entry_slide(entry: Entry, image_paths: list[Path],
                      placements_text: list[str] | None = None) -> Presentation:
    """A standalone one-scientist presentation (the Download Slide file)."""
    prs = Presentation()
    add_entry_slides(prs, entry, image_paths, placements_text)
    return prs


def build_deck(chapters: list[DeckChapter], *, deck_title: str | None = None) -> Presentation:
    """A chapter-ordered deck with a divider slide per chapter.

    ``chapters`` must already be in display order; a scientist placed in
    two chapters appears under both (Emmy Noether under Ch9 and Ch10 is
    by design), but never twice under one chapter.
    """
    prs = Presentation()
    if deck_title:
        _add_title_slide(prs, deck_title)
    for deck_chapter in sorted(chapters, key=lambda c: c.chapter):
        divider_title = f"Chapter {deck_chapter.chapter}"
        if deck_chapter.topics:
            divider_title += f": {deck_chapter.topics}"
        _add_title_slide(prs, divider_title)
        seen: set[str] = set()
        for entry, image_paths in deck_chapter.entries:
            if entry.name in seen:
                continue
            seen.add(entry.name)
            add_entry_slides(prs, entry, image_paths)
    return prs
